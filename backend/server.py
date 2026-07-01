"""
Deck Agent — backend
====================================================================
A small FastAPI app that powers the whole product:

  1. Upload a deck (PDF / PPTX / PPT)
  2. Convert to PDF if needed (LibreOffice headless), then render every
     page to an image + extract its text (PyMuPDF)
  3. Generate spoken narration grounded in the extracted text (Anthropic)
  4. Publish -> stored server-side and served as a shareable full-screen
     viewer link that anyone can open
  5. Answer prospect questions grounded ONLY in the deck, escalating to the
     rep when it can't answer confidently

See DOCS.md for the full architecture. Every endpoint lives here and is
grouped by the numbered sections below.
"""

from __future__ import annotations

import io
import json
import os
import shutil
import subprocess
import uuid
from datetime import datetime, timezone
from pathlib import Path

import fitz  # PyMuPDF
from dotenv import load_dotenv
from fastapi import FastAPI, File, HTTPException, Request, UploadFile
from fastapi.responses import FileResponse, JSONResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel

from llm import call_llm, provider_status

# ---------------------------------------------------------------------------
# 0. Config & paths
# ---------------------------------------------------------------------------
BASE_DIR = Path(__file__).resolve().parent
PROJECT_DIR = BASE_DIR.parent
FRONTEND_DIR = PROJECT_DIR / "frontend"
STORAGE_DIR = BASE_DIR / "storage"
DECKS_DIR = STORAGE_DIR / "decks"
DECKS_DIR.mkdir(parents=True, exist_ok=True)

load_dotenv(BASE_DIR / ".env")

# LLM provider config lives in llm.py (switchable via LLM_PROVIDER in .env).
SOFFICE_PATH = os.getenv("SOFFICE_PATH", "").strip()

PAGE_DPI = 140          # full page render quality
THUMB_SCALE = 0.28      # thumbnail size relative to page
MAX_TEXT_PER_PAGE = 1200  # chars of page text sent to the model per page

app = FastAPI(title="Deck Agent")


# ---------------------------------------------------------------------------
# 1. Small helpers
# ---------------------------------------------------------------------------
def now_iso() -> str:
    return datetime.now(timezone.utc).isoformat()


def deck_dir(deck_id: str) -> Path:
    return DECKS_DIR / deck_id


def deck_json_path(deck_id: str) -> Path:
    return deck_dir(deck_id) / "deck.json"


def load_deck(deck_id: str) -> dict:
    p = deck_json_path(deck_id)
    if not p.exists():
        raise HTTPException(status_code=404, detail="Deck not found")
    return json.loads(p.read_text(encoding="utf-8"))


def save_deck(deck: dict) -> None:
    deck_json_path(deck["id"]).write_text(
        json.dumps(deck, ensure_ascii=False, indent=2), encoding="utf-8"
    )


def find_soffice() -> str | None:
    """Locate the LibreOffice binary used to convert PPT/PPTX -> PDF."""
    if SOFFICE_PATH and Path(SOFFICE_PATH).exists():
        return SOFFICE_PATH
    for name in ("soffice", "libreoffice"):
        found = shutil.which(name)
        if found:
            return found
    candidates = [
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/usr/bin/soffice",
        "/usr/bin/libreoffice",
    ]
    for c in candidates:
        if Path(c).exists():
            return c
    return None


# ---------------------------------------------------------------------------
# 2. Upload -> convert -> render pages -> extract text
# ---------------------------------------------------------------------------
def convert_to_pdf(src: Path, out_dir: Path) -> Path | None:
    """Convert PPT/PPTX to PDF with LibreOffice. Returns None if unavailable."""
    soffice = find_soffice()
    if not soffice:
        return None
    subprocess.run(
        [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(out_dir), str(src)],
        check=True,
        capture_output=True,
        timeout=180,
    )
    pdf = out_dir / (src.stem + ".pdf")
    return pdf if pdf.exists() else None


def render_pdf(pdf_path: Path, ddir: Path) -> list[dict]:
    """Render each PDF page to a full image + thumbnail and extract its text."""
    pages_dir = ddir / "pages"
    thumbs_dir = ddir / "thumbs"
    pages_dir.mkdir(exist_ok=True)
    thumbs_dir.mkdir(exist_ok=True)

    doc = fitz.open(pdf_path)
    pages = []
    for i, page in enumerate(doc):
        n = i + 1
        text = (page.get_text() or "").strip()

        pix = page.get_pixmap(dpi=PAGE_DPI)
        pix.save(pages_dir / f"{n}.png")

        thumb = page.get_pixmap(matrix=fitz.Matrix(THUMB_SCALE, THUMB_SCALE))
        thumb.save(thumbs_dir / f"{n}.png")

        pages.append({"index": n, "text": text, "narration": ""})
    doc.close()
    return pages


def extract_pptx_text(src: Path, ddir: Path) -> list[dict]:
    """Fallback when LibreOffice is missing: text-only, no page images."""
    from pptx import Presentation

    prs = Presentation(str(src))
    pages = []
    for i, slide in enumerate(prs.slides):
        chunks = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        chunks.append(line)
        pages.append({"index": i + 1, "text": "\n".join(chunks), "narration": ""})
    return pages


@app.post("/api/decks")
async def create_deck(file: UploadFile = File(...)):
    ext = Path(file.filename or "").suffix.lower()
    if ext not in (".pdf", ".pptx", ".ppt"):
        raise HTTPException(status_code=400, detail="Upload a .pdf, .pptx, or .ppt file")

    deck_id = uuid.uuid4().hex[:10]
    ddir = deck_dir(deck_id)
    ddir.mkdir(parents=True, exist_ok=True)

    src = ddir / f"source{ext}"
    src.write_bytes(await file.read())

    rendered = True
    warning = None

    if ext == ".pdf":
        pdf_path = src
    else:
        pdf_path = convert_to_pdf(src, ddir)

    if pdf_path is not None:
        pages = render_pdf(pdf_path, ddir)
    elif ext == ".pptx":
        pages = extract_pptx_text(src, ddir)
        rendered = False
        warning = ("LibreOffice not found — extracted text only. Install LibreOffice "
                   "(or set SOFFICE_PATH) to get rendered page images and thumbnails.")
    else:  # .ppt with no converter
        shutil.rmtree(ddir, ignore_errors=True)
        raise HTTPException(
            status_code=422,
            detail="Legacy .ppt needs LibreOffice to be read. Install it (or set "
                   "SOFFICE_PATH), or convert the file to PDF/PPTX and re-upload.",
        )

    deck = {
        "id": deck_id,
        "title": Path(file.filename).stem if file.filename else "Untitled deck",
        "source_ext": ext,
        "rendered": rendered,
        "created_at": now_iso(),
        "published": False,
        "voice": None,
        "rep_name": "the rep",
        "pages": pages,
        "questions": [],
    }
    save_deck(deck)
    return JSONResponse({"deck": public_deck(deck, str_request_base=None), "warning": warning})


def public_deck(deck: dict, str_request_base: str | None) -> dict:
    """Shape a deck for the frontend, adding image/thumb URLs."""
    out = {k: deck[k] for k in ("id", "title", "rendered", "published", "voice", "rep_name", "created_at")}
    out["pages"] = []
    for p in deck["pages"]:
        n = p["index"]
        img = f"/files/decks/{deck['id']}/pages/{n}.png" if deck["rendered"] else None
        thumb = f"/files/decks/{deck['id']}/thumbs/{n}.png" if deck["rendered"] else None
        out["pages"].append({
            "index": n,
            "text": p.get("text", ""),
            "narration": p.get("narration", ""),
            "image": img,
            "thumb": thumb,
        })
    return out


# ---------------------------------------------------------------------------
# 3. Model proxy: narration + grounded Q&A (key stays server-side, see llm.py)
# ---------------------------------------------------------------------------
@app.post("/api/decks/{deck_id}/narrate")
async def narrate(deck_id: str):
    deck = load_deck(deck_id)
    system = (
        "You write spoken narration scripts for B2B sales pitch decks, read aloud by "
        "text-to-speech, one segment per slide, playing in sync as each slide is shown. "
        "Write exactly one spoken line per slide (max 32 words), building a coherent arc "
        "across the slides in order. Tone: confident, plain, like a skilled human rep "
        "talking to a busy prospect. No hype adjectives, no exclamation marks, no emoji. "
        "Return ONLY a raw JSON array of strings, same length and order as the input "
        "slides, nothing else — no markdown, no code fences, no commentary."
    )
    slides_payload = [
        {"slide": p["index"], "text": (p.get("text") or "")[:MAX_TEXT_PER_PAGE]}
        for p in deck["pages"]
    ]
    raw = await call_llm(system, json.dumps(slides_payload), max_tokens=1600)
    try:
        lines = json.loads(raw)
        assert isinstance(lines, list) and len(lines) == len(deck["pages"])
    except Exception:
        raise HTTPException(status_code=502, detail="Model returned a mismatched script — try again.")
    for i, line in enumerate(lines):
        deck["pages"][i]["narration"] = str(line)
    save_deck(deck)
    return {"narration": [p["narration"] for p in deck["pages"]]}


class AskBody(BaseModel):
    question: str


@app.post("/api/decks/{deck_id}/ask")
async def ask(deck_id: str, body: AskBody):
    deck = load_deck(deck_id)
    question = (body.question or "").strip()
    if not question:
        raise HTTPException(status_code=400, detail="Empty question")

    rep = deck.get("rep_name") or "the rep"
    deck_text = "\n\n".join(
        f"Slide {p['index']}:\n{(p.get('text') or '').strip()}" for p in deck["pages"]
    )
    system = (
        "You answer a prospect's question about a sales deck they just viewed. You may "
        "ONLY use the deck content below — never use outside knowledge, never invent "
        "numbers, pricing, dates, security/compliance claims, or commitments not present "
        "in the text.\n\nDECK CONTENT:\n" + deck_text +
        '\n\nRespond with ONLY raw JSON, no markdown, in exactly this shape:\n'
        '{"escalate": boolean, "answer": "string", "slide_ref": number or null, '
        '"confidence": number between 0 and 1}\n\nRules:\n'
        "- If the deck clearly answers the question: escalate=false, confidence>=0.7, "
        "slide_ref=the 1-indexed slide, answer=a plain 1-2 sentence answer in a human, "
        "non-salesy voice.\n"
        "- If the question needs anything not in the deck (exact contract terms, custom "
        "pricing, security/compliance detail, unstated timelines, or anything you're not "
        "confident about): escalate=true, confidence<0.5, slide_ref=null, and "
        f'answer="Good question — let me get {rep} to answer that directly for you."'
    )
    raw = await call_llm(system, question, max_tokens=400)
    try:
        parsed = json.loads(raw)
    except Exception:
        parsed = {"escalate": True, "slide_ref": None, "confidence": 0.0,
                  "answer": f"Good question — let me get {rep} to answer that directly for you."}

    record = {
        "question": question,
        "escalate": bool(parsed.get("escalate")),
        "answer": parsed.get("answer", ""),
        "slide_ref": parsed.get("slide_ref"),
        "confidence": parsed.get("confidence"),
        "ts": now_iso(),
    }
    deck.setdefault("questions", []).append(record)
    save_deck(deck)
    return parsed


# ---------------------------------------------------------------------------
# 4. Publish + fetch + rep inbox
# ---------------------------------------------------------------------------
class PublishBody(BaseModel):
    narration: list[str] | None = None
    voice: dict | None = None
    rep_name: str | None = None
    title: str | None = None


@app.post("/api/decks/{deck_id}/publish")
async def publish(deck_id: str, body: PublishBody, request: Request):
    deck = load_deck(deck_id)
    if body.narration is not None:
        if len(body.narration) != len(deck["pages"]):
            raise HTTPException(status_code=400, detail="Narration length mismatch")
        for i, line in enumerate(body.narration):
            deck["pages"][i]["narration"] = line
    if body.voice is not None:
        deck["voice"] = body.voice
    if body.rep_name:
        deck["rep_name"] = body.rep_name
    if body.title:
        deck["title"] = body.title
    deck["published"] = True
    deck["published_at"] = now_iso()
    save_deck(deck)

    base = str(request.base_url).rstrip("/")
    return {"url": f"{base}/d/{deck_id}", "path": f"/d/{deck_id}"}


@app.get("/api/decks/{deck_id}")
async def get_deck(deck_id: str):
    return public_deck(load_deck(deck_id), str_request_base=None)


@app.get("/api/decks/{deck_id}/questions")
async def get_questions(deck_id: str):
    deck = load_deck(deck_id)
    return {"rep_name": deck.get("rep_name"), "questions": deck.get("questions", [])}


# ---------------------------------------------------------------------------
# 5. Static files & page routes (mounted LAST so /api wins)
# ---------------------------------------------------------------------------
app.mount("/files", StaticFiles(directory=STORAGE_DIR), name="files")


@app.get("/d/{deck_id}")
async def viewer_page(deck_id: str):
    # The viewer reads the deck id from its own URL and calls /api/decks/{id}.
    return FileResponse(FRONTEND_DIR / "viewer.html")


@app.get("/health")
async def health():
    return {"ok": True, "llm": provider_status(), "soffice": find_soffice()}


# Studio (index.html) + assets at the root. html=True serves index.html for "/".
app.mount("/", StaticFiles(directory=FRONTEND_DIR, html=True), name="frontend")
